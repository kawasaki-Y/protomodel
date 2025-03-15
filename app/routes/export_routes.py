from flask import Blueprint, render_template, request, jsonify, Response, current_app
import tempfile
import os
from weasyprint import HTML, CSS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io
import json
from app.models.models import BusinessPlan, PlanItem, CashFlowPlan, CashFlowItem
from datetime import datetime
from pptx.dml.chart import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from io import BytesIO
from flask import make_response
from flask_login import login_required, current_user

export_bp = Blueprint('export', __name__)

@export_bp.route('/api/export/pdf', methods=['POST'])
@login_required
def export_pdf():
    """事業計画データをPDFに出力する"""
    try:
        data = request.json
        plan_id = data.get('plan_id')
        include_cash_flow = data.get('include_cash_flow', True)
        
        # 事業計画データを取得
        plan = BusinessPlan.query.get(plan_id)
        if not plan:
            return jsonify({'error': '事業計画が見つかりません'}), 404
        
        # 事業計画項目のデータを取得（階層含む）
        items = PlanItem.query.filter_by(business_plan_id=plan_id).all()
        
        # 資金繰り計画を取得
        cash_flow_items = None
        if include_cash_flow:
            # 関連する資金繰り計画を取得
            cash_flow_plan = CashFlowPlan.query.filter_by(business_plan_id=plan_id).first()
            if cash_flow_plan:
                cash_flow_items = CashFlowItem.query.filter_by(cash_flow_plan_id=cash_flow_plan.id).all()
        
        # 現在の日付を取得
        current_date = datetime.now().strftime('%Y年%m月%d日')
        current_year = datetime.now().year
        
        # HTMLレポートの生成
        html = render_template(
            'export/pdf_report.html',
            plan=plan,
            items=items,
            cash_flow_items=cash_flow_items,
            current_date=current_date,
            current_year=current_year
        )
        
        # HTML -> PDF変換
        pdf = HTML(string=html, base_url=request.url_root).write_pdf()
        
        # レスポンスを返す
        response = make_response(pdf)
        response.headers['Content-Type'] = 'application/pdf'
        response.headers['Content-Disposition'] = f'attachment; filename=business_plan_{plan.id}.pdf'
        
        return response
    except Exception as e:
        current_app.logger.error(f"PDF生成エラー: {str(e)}")
        return jsonify({'error': f'PDF生成中にエラーが発生しました: {str(e)}'}), 500

@export_bp.route('/api/export/ppt', methods=['POST'])
@login_required
def export_ppt():
    """事業計画データをPowerPointに出力する"""
    try:
        data = request.json
        plan_id = data.get('plan_id')
        include_chart = data.get('include_chart', True)
        
        # 事業計画データを取得
        plan = BusinessPlan.query.get(plan_id)
        if not plan:
            return jsonify({'error': '事業計画が見つかりません'}), 404
        
        # 事業計画項目のデータを取得
        items = PlanItem.query.filter_by(business_plan_id=plan_id).all()
        
        # PowerPointプレゼンテーションの作成
        prs = Presentation()
        
        # タイトルスライド
        slide_layout = prs.slide_layouts[0]  # タイトル用レイアウト
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{plan.name}"
        subtitle.text = f"{plan.fiscal_year}年度 事業計画 \n作成日: {datetime.now().strftime('%Y年%m月%d日')}"
        
        # 概要スライド
        slide_layout = prs.slide_layouts[1]  # タイトルと内容のレイアウト
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "事業計画の概要"
        
        # 売上と利益のデータを取得
        sales_items = [item for item in items if '売上' in item.name]
        profit_items = [item for item in items if '利益' in item.name or '営業利益' in item.name]
        
        content_text = ""
        if sales_items:
            content_text += f"売上目標: {sum(getattr(sales_items[0], f'm{i}_amount', 0) for i in range(1, 13)):,}円\n"
        if profit_items:
            content_text += f"利益目標: {sum(getattr(profit_items[0], f'm{i}_amount', 0) for i in range(1, 13)):,}円\n"
        
        content_text += f"\n{plan.description or '本事業計画は年間の事業目標を達成するための指針となるものです。'}"
        content.text = content_text
        
        # グラフスライド（オプション）
        if include_chart and (sales_items or profit_items):
            slide_layout = prs.slide_layouts[5]  # タイトルのみのレイアウト
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "月次の売上・利益計画"
            
            # グラフデータの準備
            chart_data = CategoryChartData()
            chart_data.categories = ['1月', '2月', '3月', '4月', '5月', '6月', 
                                    '7月', '8月', '9月', '10月', '11月', '12月']
            
            # 売上データ
            if sales_items:
                sales_data = [getattr(sales_items[0], f'm{i}_amount', 0) for i in range(1, 13)]
                chart_data.add_series('売上', sales_data)
            
            # 利益データ
            if profit_items:
                profit_data = [getattr(profit_items[0], f'm{i}_amount', 0) for i in range(1, 13)]
                chart_data.add_series('利益', profit_data)
            
            # グラフの追加
            x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        
        # PowerPointを一時ファイルに保存
        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        
        # レスポンスを返す
        response = make_response(pptx_stream.getvalue())
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        response.headers['Content-Disposition'] = f'attachment; filename=business_plan_{plan.id}.pptx'
        
        return response
    except Exception as e:
        current_app.logger.error(f"PPT生成エラー: {str(e)}")
        return jsonify({'error': f'PowerPoint生成中にエラーが発生しました: {str(e)}'}), 500 